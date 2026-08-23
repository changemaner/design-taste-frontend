# -*- coding: utf-8 -*-
"""
deck 导出 PPTX（P2 部署/导出闭环）
================================
每套 deck 导出为 16:9 PowerPoint：逐页截 1920×1080 全幅图嵌入幻灯片，
保留全部视觉与字体渲染结果（与浏览器所见一致）。

用法：
  python showcase/export-pptx.py <slug>            # showcase 样品 → showcase/exports/<slug>.pptx
  python showcase/export-pptx.py <deck-dir>        # 用户自己的 deck 目录（包外任意含
                                                    # index.html 的目录）→ <deck-dir>/<name>.pptx
  python showcase/export-pptx.py <x> --out dir

deck 经 file:// 直读，无需起本地 HTTP 服务；deck 引用的 ../fonts/fonts.css 等
相对资源按目录结构解析（字体库从 showcase/fonts/ 拷一份到 deck 同级即可，
与部署暂存同结构）。
前置：playwright + Chrome + python-pptx（pip install python-pptx）。
"""
import asyncio
import sys
import tempfile
from pathlib import Path

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
SHOWCASE = ROOT / "showcase"
def find_chrome():
    """Chrome 定位：DESIGN_CHROME / CHROME_PATH 环境变量 → 各平台常见安装路径 → None。
    返回 None 时 playwright 回退自带 chromium（需 playwright install chromium）。"""
    import os
    for var in ("DESIGN_CHROME", "CHROME_PATH"):
        p = os.environ.get(var)
        if p and Path(p).exists():
            return p
    for c in (
        r"C:\Program Files\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
        str(Path.home() / "AppData" / "Local" / "Google" / "Chrome" / "Application" / "chrome.exe"),
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/usr/bin/google-chrome",
        "/usr/bin/chromium",
        "/usr/bin/chromium-browser",
        "/snap/bin/chromium",
    ):
        if Path(c).exists():
            return c
    return None

SLIDE_W_IN = 13.3333  # 16:9
SLIDE_H_IN = 7.5


def resolve_deck(arg: str) -> Path:
    """slug → showcase/<slug>/；否则当作 deck 目录路径（用户自己生成的 HTML）。"""
    slug_dir = SHOWCASE / arg
    if (slug_dir / "index.html").exists():
        return slug_dir
    p = Path(arg)
    if (p / "index.html").exists():
        return p.resolve()
    raise SystemExit(f"[!!] 找不到 deck：{arg}（既不是 showcase/<slug>，也不是含 index.html 的目录）")


async def export_pptx(deck: Path, out: Path):
    url = (deck / "index.html").as_uri()
    out.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory() as tmp:
        tmp = Path(tmp)
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True, executable_path=find_chrome())
            page = await browser.new_page(viewport={"width": 1920, "height": 1080})
            await page.goto(url)
            n_slides = await page.evaluate("document.querySelectorAll('.slide').length")
            shots = []
            for i in range(n_slides):
                await page.goto(f"{url}?slide={i}")
                await page.evaluate("document.fonts.ready")
                await page.wait_for_timeout(1500)
                shot = tmp / f"slide{i}.png"
                await page.screenshot(path=str(shot))
                shots.append(shot)
            await browser.close()

        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W_IN)
        prs.slide_height = Inches(SLIDE_H_IN)
        blank = prs.slide_layouts[6]  # blank layout
        for shot in shots:
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(str(shot), 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
        prs.save(out)
    try:
        shown = out.relative_to(ROOT)
    except ValueError:
        shown = out.resolve()
    print(f"[ok] {deck.name} -> {shown}（{n_slides} 页）")


def main():
    if len(sys.argv) < 2:
        raise SystemExit(__doc__)
    deck = resolve_deck(sys.argv[1])
    if "--out" in sys.argv:
        out_dir = Path(sys.argv[sys.argv.index("--out") + 1])
    elif deck.parent == SHOWCASE:
        out_dir = SHOWCASE / "exports"  # showcase 样品：维持 exports/ 惯例
    else:
        out_dir = deck.parent  # 包外 deck：产物放 deck 旁边
    asyncio.run(export_pptx(deck, out_dir / f"{deck.name}.pptx"))


if __name__ == "__main__":
    main()
